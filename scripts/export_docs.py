from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.shapes import Circle, Drawing, Line, Rect, String
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
DOCS = ASSETS / "docs"
IMAGES = ASSETS / "images"


NAVY = colors.HexColor("#1E2A3A")
SAND = colors.HexColor("#C2A57A")
SAND_SOFT = colors.HexColor("#EBDCC6")
TERRACOTTA = colors.HexColor("#B96939")
OLIVE = colors.HexColor("#617147")
CREAM = colors.HexColor("#FFFAF3")
INK = colors.HexColor("#1F2A36")
MUTED = colors.HexColor("#65707A")
LINE = colors.HexColor("#D8CFC2")
WHITE = colors.white


YEAR_LABELS = ["Year 1", "Year 2", "Year 3"]
AUDIENCE = [50000, 250000, 1000000]
COLLECTIONS = [4, 8, 12]
TOTAL_UNITS = [10000, 40000, 96000]
AVERAGE_SELLING_PRICE = [65, 70, 72]
PRODUCT_REVENUE = [650000, 2800000, 6912000]
MEMBERSHIP_COUNT = [1000, 5000, 15000]
MEMBERSHIP_CONVERSION = [0.02, 0.02, 0.015]
MEMBERSHIP_PRICE = 8
MEMBERSHIP_REVENUE = [96000, 480000, 1440000]
PARTNERSHIP_REVENUE = [80000, 320000, 1000000]
REVENUE = [product + membership + partnership for product, membership, partnership in zip(PRODUCT_REVENUE, MEMBERSHIP_REVENUE, PARTNERSHIP_REVENUE)]
COGS = [280000, 1120000, 2688000]
TRAVEL_AND_SOURCING = [60000, 85000, 120000]
TEAM_COST = [150000, 450000, 900000]
MARKETING_AND_GROWTH = [82600, 432000, 1028720]
GENERAL_OPS = [75000, 180000, 350000]
OPERATING_COSTS = [travel + team + marketing + ops for travel, team, marketing, ops in zip(TRAVEL_AND_SOURCING, TEAM_COST, MARKETING_AND_GROWTH, GENERAL_OPS)]
GROSS_PROFIT = [revenue - cogs for revenue, cogs in zip(REVENUE, COGS)]
OPERATING_PROFIT = [gross_profit - operating_cost for gross_profit, operating_cost in zip(GROSS_PROFIT, OPERATING_COSTS)]
GROSS_MARGIN = [gross_profit / revenue for gross_profit, revenue in zip(GROSS_PROFIT, REVENUE)]
OPERATING_MARGIN = [operating_profit / revenue for operating_profit, revenue in zip(OPERATING_PROFIT, REVENUE)]


@dataclass(frozen=True)
class ProductCard:
    name: str
    note: str


@dataclass(frozen=True)
class RegionCard:
    name: str
    note: str


@dataclass(frozen=True)
class SkuModel:
    name: str
    role: str
    unit_mix: float
    asp: tuple[float, float, float]
    direct_cost: tuple[float, float, float]


PRODUCTS = [
    ProductCard("Indigo Linen Throw", "Hero product with premium texture, gifting appeal, and strong editorial presence."),
    ProductCard("Indigo Scarf", "Portable entry point that supports travel, styling, and repeat purchase."),
    ProductCard("Indigo Pillow", "Home category anchor that translates textile story into everyday living."),
    ProductCard("Indigo Textile Set", "Study set that turns process, swatch, and provenance into a collectible object."),
]

REGIONS = [
    RegionCard("Japan", "Indigo process, heritage finishing, and the Kyoto Indigo narrative foundation."),
    RegionCard("India", "Block print, handloom depth, flexible development capacity, and future assortment breadth."),
    RegionCard("Turkey", "Cotton scale, towels, throws, and production fluency for broader textile programs."),
    RegionCard("Italy", "Luxury finishing, linen credibility, and eventual prestige textile manufacturing."),
    RegionCard("Peru", "Future expansion into alpaca and rich fiber stories for premium soft goods."),
]

SKU_MODELS = [
    SkuModel("Indigo Linen Throw", "Hero textile anchor with the strongest visual and gifting appeal.", 0.18, (118, 128, 132), (52, 51, 52)),
    SkuModel("Indigo Scarf", "Portable entry price point that keeps the collection accessible and repeatable.", 0.34, (52, 56, 58), (22, 22, 22)),
    SkuModel("Indigo Pillow", "Home category bridge that brings the collection into everyday living.", 0.24, (69, 73, 75), (30, 30, 30)),
    SkuModel("Indigo Textile Set", "Collector-oriented accessory product with strong storytelling value.", 0.24, (40, 44, 45), (17, 17, 17)),
]

LAUNCH_CADENCE = [
    ("Year 1", "4 launches", "Kyoto Indigo, Amalfi Linen, Marrakech Weave, Anatolia Cotton", "Build recognition around one destination story at a time."),
    ("Year 2", "8 launches", "Repeat winners plus Jaipur, Tuscany, Oaxaca, and Denizli textile programs", "Broaden assortment carefully while preserving the destination-led model."),
    ("Year 3", "12 launches", "Monthly cadence across hero destinations and early private label extensions", "Move from proof into category authority and signature lines."),
]

CAPITAL_PLAN = [
    ("Product development & samples", 12000, 85000, "Sampling, refinement, and early private label design work."),
    ("Inventory", 22000, 180000, "Opening buys, reorders, and broader assortment depth."),
    ("Travel & sourcing", 10000, 90000, "Supplier visits, quality control, and regional discovery work."),
    ("Brand & content production", 12000, 100000, "Editorial, video, product storytelling, and launch assets."),
    ("Commerce & fulfillment systems", 8000, 45000, "Platform, fulfillment setup, and operational infrastructure."),
    ("Team build-out", 6000, 250000, "Design, operations, growth, and supplier management capacity."),
]


def money_short(value: float) -> str:
    if value >= 1_000_000:
        return f"${value / 1_000_000:.1f}M"
    if value >= 1_000:
        return f"${value / 1_000:.0f}K"
    return f"${value:,.0f}"


def ensure_dirs() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)


def image_dimensions(path: Path) -> tuple[int, int]:
    with PILImage.open(path) as img:
        return img.size


def fitted_image(path: Path, max_width: float, max_height: float, h_align: str = "CENTER") -> Image:
    width_px, height_px = ImageReader(str(path)).getSize()
    scale = min(max_width / width_px, max_height / height_px)
    flowable = Image(str(path), width=width_px * scale, height=height_px * scale)
    flowable.hAlign = h_align
    return flowable


def build_styles():
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="CoverKicker",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=10.5,
            leading=12,
            textColor=SAND,
            spaceAfter=10,
            alignment=TA_CENTER,
        )
    )
    styles.add(
        ParagraphStyle(
            name="CoverTitle",
            parent=styles["Title"],
            fontName="Times-Bold",
            fontSize=30,
            leading=34,
            textColor=NAVY,
            alignment=TA_CENTER,
            spaceAfter=14,
        )
    )
    styles.add(
        ParagraphStyle(
            name="CoverBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=12.5,
            leading=18,
            textColor=MUTED,
            alignment=TA_CENTER,
            spaceAfter=14,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SectionKicker",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=9.5,
            leading=11,
            textColor=TERRACOTTA,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SectionTitle",
            parent=styles["Heading1"],
            fontName="Times-Bold",
            fontSize=22,
            leading=26,
            textColor=NAVY,
            spaceAfter=10,
        )
    )
    styles.add(
        ParagraphStyle(
            name="BodyCopy",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.7,
            leading=15,
            textColor=INK,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="BodyCopyMuted",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.2,
            leading=14,
            textColor=MUTED,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="BulletCopy",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10.3,
            leading=14,
            leftIndent=14,
            bulletIndent=2,
            textColor=INK,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="CardTitle",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=11.3,
            leading=13,
            textColor=NAVY,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="CardBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=9.7,
            leading=13.2,
            textColor=INK,
            spaceAfter=2,
        )
    )
    styles.add(
        ParagraphStyle(
            name="StatValue",
            parent=styles["BodyText"],
            fontName="Times-Bold",
            fontSize=21,
            leading=23,
            textColor=NAVY,
            alignment=TA_CENTER,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="StatLabel",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=8.8,
            leading=10.5,
            textColor=MUTED,
            alignment=TA_CENTER,
            spaceAfter=2,
        )
    )
    styles.add(
        ParagraphStyle(
            name="MiniLabel",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=8.5,
            leading=10,
            textColor=MUTED,
            alignment=TA_CENTER,
        )
    )
    return styles


def card(content, *, background=WHITE, border_color=LINE, padding=12, valign="TOP") -> Table:
    table = Table([[content]])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), background),
                ("BOX", (0, 0), (-1, -1), 0.8, border_color),
                ("INNERPADDING", (0, 0), (-1, -1), padding),
                ("VALIGN", (0, 0), (-1, -1), valign),
                ("ROUNDEDCORNERS", (0, 0), (-1, -1), 10),
            ]
        )
    )
    return table


def stat_card(value: str, label: str, styles) -> Table:
    return card(
        [
            Paragraph(value, styles["StatValue"]),
            Paragraph(label, styles["StatLabel"]),
        ],
        background=WHITE,
        padding=10,
        valign="MIDDLE",
    )


def narrative_card(title: str, body: str, styles, *, background=WHITE) -> Table:
    return card(
        [
            Paragraph(title, styles["CardTitle"]),
            Paragraph(body, styles["CardBody"]),
        ],
        background=background,
        padding=12,
    )


def bullets(items: list[str], styles) -> list[Paragraph]:
    return [Paragraph(item, styles["BulletCopy"], bulletText="•") for item in items]


def section_heading(kicker: str, title: str, body: str, styles) -> list:
    return [
        Paragraph(kicker, styles["SectionKicker"]),
        Paragraph(title, styles["SectionTitle"]),
        Paragraph(body, styles["BodyCopy"]),
    ]


def revenue_chart(width: float = 6.7 * inch, height: float = 2.8 * inch) -> Drawing:
    drawing = Drawing(width, height)
    drawing.add(Rect(0, 0, width, height, fillColor=WHITE, strokeColor=LINE, rx=12, ry=12))

    chart = VerticalBarChart()
    chart.x = 34
    chart.y = 34
    chart.width = width - 68
    chart.height = height - 74
    chart.data = [[value / 1_000_000 for value in REVENUE], [value / 1_000_000 for value in GROSS_PROFIT]]
    chart.valueAxis.valueMin = 0
    chart.valueAxis.valueMax = 10
    chart.valueAxis.valueStep = 2
    chart.valueAxis.labels.fontName = "Helvetica"
    chart.valueAxis.labels.fontSize = 7
    chart.valueAxis.labels.fillColor = MUTED
    chart.categoryAxis.categoryNames = YEAR_LABELS
    chart.categoryAxis.labels.fontName = "Helvetica-Bold"
    chart.categoryAxis.labels.fontSize = 8
    chart.categoryAxis.labels.fillColor = NAVY
    chart.categoryAxis.strokeColor = colors.transparent
    chart.valueAxis.strokeColor = colors.transparent
    chart.bars[0].fillColor = NAVY
    chart.bars[0].strokeColor = NAVY
    chart.bars[1].fillColor = TERRACOTTA
    chart.bars[1].strokeColor = TERRACOTTA
    chart.barWidth = 18
    chart.groupSpacing = 16
    chart.barSpacing = 6
    chart.categoryAxis.tickDown = 0
    drawing.add(chart)

    drawing.add(String(18, height - 18, "Revenue and gross profit ($M)", fontName="Helvetica-Bold", fontSize=10, fillColor=NAVY))
    drawing.add(Rect(width - 150, height - 22, 8, 8, fillColor=NAVY, strokeColor=NAVY))
    drawing.add(String(width - 136, height - 20, "Revenue", fontName="Helvetica", fontSize=8, fillColor=MUTED))
    drawing.add(Rect(width - 82, height - 22, 8, 8, fillColor=TERRACOTTA, strokeColor=TERRACOTTA))
    drawing.add(String(width - 68, height - 20, "Gross profit", fontName="Helvetica", fontSize=8, fillColor=MUTED))
    return drawing


def roadmap_chart(width: float = 6.7 * inch, height: float = 2.05 * inch) -> Drawing:
    drawing = Drawing(width, height)
    drawing.add(Rect(0, 0, width, height, fillColor=WHITE, strokeColor=LINE, rx=12, ry=12))
    y = height / 2
    x_positions = [width * 0.17, width * 0.5, width * 0.83]
    drawing.add(Line(46, y, width - 46, y, strokeColor=SAND, strokeWidth=3))
    phases = [
        ("Discovery", "Years 1-2", "Destination-led launches and clear brand imagery"),
        ("Authority", "Years 3-4", "Repeat collections, supplier depth, private label"),
        ("Expansion", "Years 5+", "Home textiles, prestige retail, licensing"),
    ]
    for x, (title, period, note) in zip(x_positions, phases):
        drawing.add(Circle(x, y, 13, fillColor=NAVY, strokeColor=NAVY))
        drawing.add(Circle(x, y, 5, fillColor=SAND, strokeColor=SAND))
        drawing.add(String(x - 34, y + 22, title, fontName="Helvetica-Bold", fontSize=10, fillColor=NAVY))
        drawing.add(String(x - 34, y + 10, period, fontName="Helvetica", fontSize=8, fillColor=TERRACOTTA))
        drawing.add(String(x - 64, y - 30, note, fontName="Helvetica", fontSize=7.3, fillColor=MUTED))
    return drawing


def financial_table(styles) -> Table:
    rows = [
        ["Metric", "Year 1", "Year 2", "Year 3"],
        ["Total revenue", money_short(REVENUE[0]), money_short(REVENUE[1]), money_short(REVENUE[2])],
        ["Gross profit", money_short(GROSS_PROFIT[0]), money_short(GROSS_PROFIT[1]), money_short(GROSS_PROFIT[2])],
        ["Operating profit", money_short(OPERATING_PROFIT[0]), money_short(OPERATING_PROFIT[1]), money_short(OPERATING_PROFIT[2])],
        ["Collections launched", "4", "8", "12"],
    ]
    table = Table(rows, colWidths=[2.0 * inch, 1.4 * inch, 1.4 * inch, 1.4 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), NAVY),
                ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BACKGROUND", (0, 1), (-1, -1), WHITE),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, colors.HexColor("#FBF6EF")]),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("TEXTCOLOR", (0, 1), (-1, -1), INK),
                ("GRID", (0, 0), (-1, -1), 0.6, LINE),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                ("ALIGN", (1, 1), (-1, -1), "CENTER"),
            ]
        )
    )
    return table


def unit_economics_table() -> Table:
    rows = [
        ["Unit economics", "Value"],
        ["Average selling price range", f"${AVERAGE_SELLING_PRICE[0]}-${AVERAGE_SELLING_PRICE[-1]}"],
        ["Average direct cost / unit", "$28"],
        ["Gross margin range", "55%-65%"],
        ["Membership price", f"${MEMBERSHIP_PRICE} / month"],
        ["Initial capital range", "$35K-$70K"],
    ]
    table = Table(rows, colWidths=[2.5 * inch, 1.9 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), SAND_SOFT),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("TEXTCOLOR", (0, 0), (-1, 0), NAVY),
                ("GRID", (0, 0), (-1, -1), 0.6, LINE),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, colors.HexColor("#FBF6EF")]),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    return table


def revenue_mix_table() -> Table:
    rows = [
        ["Revenue stream", "Year 1", "Year 2", "Year 3"],
        ["Product revenue", money_short(PRODUCT_REVENUE[0]), money_short(PRODUCT_REVENUE[1]), money_short(PRODUCT_REVENUE[2])],
        ["Membership revenue", money_short(MEMBERSHIP_REVENUE[0]), money_short(MEMBERSHIP_REVENUE[1]), money_short(MEMBERSHIP_REVENUE[2])],
        ["Partnership revenue", money_short(PARTNERSHIP_REVENUE[0]), money_short(PARTNERSHIP_REVENUE[1]), money_short(PARTNERSHIP_REVENUE[2])],
    ]
    table = Table(rows, colWidths=[2.25 * inch, 1.32 * inch, 1.32 * inch, 1.32 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), NAVY),
                ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BACKGROUND", (0, 1), (-1, -1), WHITE),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, colors.HexColor("#FBF6EF")]),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("TEXTCOLOR", (0, 1), (-1, -1), INK),
                ("GRID", (0, 0), (-1, -1), 0.6, LINE),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                ("ALIGN", (1, 1), (-1, -1), "CENTER"),
            ]
        )
    )
    return table


def capital_plan_table() -> Table:
    rows = [["Use of funds", "Initial capital", "Growth capital"]]
    for category, initial_capital, growth_capital, _ in CAPITAL_PLAN:
        rows.append([category, f"${initial_capital:,.0f}", f"${growth_capital:,.0f}"])
    rows.append([
        "Total",
        f"${sum(item[1] for item in CAPITAL_PLAN):,.0f}",
        f"${sum(item[2] for item in CAPITAL_PLAN):,.0f}",
    ])
    table = Table(rows, colWidths=[3.05 * inch, 1.6 * inch, 1.6 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), SAND_SOFT),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("TEXTCOLOR", (0, 0), (-1, 0), NAVY),
                ("ROWBACKGROUNDS", (0, 1), (-1, -2), [WHITE, colors.HexColor("#FBF6EF")]),
                ("BACKGROUND", (0, -1), (-1, -1), colors.HexColor("#EEF4E7")),
                ("FONTNAME", (0, -1), (-1, -1), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.6, LINE),
                ("TOPPADDING", (0, 0), (-1, -1), 7),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                ("ALIGN", (1, 1), (-1, -1), "CENTER"),
            ]
        )
    )
    return table


def draw_first_page(canvas, doc) -> None:
    canvas.saveState()
    width, height = letter
    canvas.setFillColor(CREAM)
    canvas.rect(0, 0, width, height, stroke=0, fill=1)
    canvas.setFillColor(NAVY)
    canvas.rect(0, height - 1.6 * inch, width, 1.6 * inch, stroke=0, fill=1)
    canvas.setFillColor(SAND)
    canvas.circle(width - 0.95 * inch, height - 0.78 * inch, 0.22 * inch, stroke=0, fill=1)
    canvas.setFillColor(TERRACOTTA)
    canvas.circle(0.95 * inch, 0.9 * inch, 0.12 * inch, stroke=0, fill=1)
    canvas.restoreState()


def draw_later_pages(canvas, doc) -> None:
    canvas.saveState()
    width, height = letter
    canvas.setStrokeColor(LINE)
    canvas.line(doc.leftMargin, height - 0.52 * inch, width - doc.rightMargin, height - 0.52 * inch)
    canvas.setFillColor(NAVY)
    canvas.setFont("Helvetica-Bold", 8)
    canvas.drawString(doc.leftMargin, height - 0.38 * inch, "ALI DANDIN PREMIUM TEXTILES")
    canvas.setFillColor(MUTED)
    canvas.setFont("Helvetica", 8)
    canvas.drawRightString(width - doc.rightMargin, 0.42 * inch, f"{canvas.getPageNumber()}")
    canvas.restoreState()


def build_pdf() -> None:
    output = DOCS / "investor_business_plan.pdf"
    styles = build_styles()
    story: list = []

    cover_left = [
        Paragraph("Premium textile brand", styles["CoverKicker"]),
        Paragraph("Ali Dandin Business Plan", styles["CoverTitle"]),
        Paragraph(
            "Ali Dandin defines a focused position in global textile commerce. The brand centers premium textiles, destination-led collections, and disciplined sourcing inside an editorial direct-to-consumer model.",
            styles["CoverBody"],
        ),
        Spacer(1, 0.1 * inch),
        Paragraph(
            "Kyoto Indigo establishes the opening collection, while the broader company scales through signature lines, private label, and selective partnerships.",
            styles["BodyCopy"],
        ),
    ]
    cover_image = fitted_image(IMAGES / "kyoto_indigo_collection.png", 2.7 * inch, 4.8 * inch)
    cover_table = Table([[cover_left, card([cover_image], background=colors.HexColor("#F7F0E5"), padding=10)]], colWidths=[4.45 * inch, 2.15 * inch])
    cover_table.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "MIDDLE"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))

    story.append(Spacer(1, 0.55 * inch))
    story.append(cover_table)
    story.append(Spacer(1, 0.22 * inch))

    metrics = Table(
        [[
            stat_card("Kyoto Indigo", "Opening collection", styles),
            stat_card("$826K", "Year 1 revenue", styles),
            stat_card(money_short(REVENUE[2]), "Year 3 revenue", styles),
            stat_card("55%-65%", "Gross margin target", styles),
        ]],
        colWidths=[1.68 * inch, 1.4 * inch, 1.4 * inch, 1.4 * inch],
    )
    metrics.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0), ("VALIGN", (0, 0), (-1, -1), "MIDDLE")]))
    story.append(metrics)
    story.append(PageBreak())

    story.extend(section_heading("Brand snapshot", "A textile-led brand with clean scale logic", "Ali Dandin combines category clarity, disciplined product architecture, and premium positioning in a way that translates well across commerce, partnerships, and future private label development.", styles))
    snapshot_cards = Table(
        [[
            narrative_card("Core thesis", "Premium textiles built around provenance, destination, and material literacy rather than generic travel commerce.", styles),
            narrative_card("Opening move", "Kyoto Indigo defines the first collection and establishes the visual and merchandising grammar.", styles),
        ], [
            narrative_card("Commercial logic", "Narrow collections support stronger pricing, cleaner storytelling, and lower early inventory risk.", styles),
            narrative_card("Expansion path", "The platform grows into repeat collections, signature lines, and selective retail or licensing relationships.", styles),
        ]],
        colWidths=[3.18 * inch, 3.18 * inch],
    )
    snapshot_cards.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    snapshot_stats = Table(
        [[
            stat_card("4", "Opening SKUs", styles),
            stat_card("12", "Collections by Year 3", styles),
            stat_card(f"{MEMBERSHIP_COUNT[-1]:,}", "Membership base", styles),
            stat_card("$35K-$70K", "Initial capital", styles),
        ]],
        colWidths=[1.56 * inch, 1.56 * inch, 1.56 * inch, 1.56 * inch],
    )
    snapshot_stats.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(snapshot_cards)
    story.append(Spacer(1, 0.16 * inch))
    story.append(snapshot_stats)
    story.append(PageBreak())

    story.extend(section_heading("Investment case", "A premium textile authority with clear commercial logic", "Ali Dandin is positioned as a specialized authority in premium textiles. The thesis is simple: provenance, material literacy, and destination-led curation create a cleaner premium brand than broad travel commerce or generic artisan marketplaces.", styles))
    thesis_cards = Table(
        [[
            narrative_card("Category edge", "Premium textiles combine tactile differentiation, repeat purchase, healthy price architecture, and eventual private label expansion.", styles),
            narrative_card("Brand edge", "Ali's background in textiles brings sourcing fluency, material judgment, and commercial credibility that most travel-led brands do not have.", styles),
            narrative_card("Growth edge", "A tightly edited collection structure makes the brand legible to customers, retailers, and future capital partners.", styles),
        ]],
        colWidths=[2.12 * inch, 2.12 * inch, 2.12 * inch],
    )
    thesis_cards.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(thesis_cards)
    story.append(Spacer(1, 0.16 * inch))
    story.append(revenue_chart())
    story.append(Spacer(1, 0.14 * inch))
    story.extend(
        bullets(
            [
                "Product revenue establishes the brand; recurring revenue and private label extend it.",
                "Destination-led collections build memory structure and keep the assortment disciplined.",
                "The value of the company grows as the brand moves from curated discovery into a repeatable textile system.",
            ],
            styles,
        )
    )
    story.append(PageBreak())

    story.extend(section_heading("Brand system", "Identity, tone, and visual authority", "The identity system includes a primary mark, monogram, passport stamp, and textile globe. The design language is restrained and editorial: midnight blue, sand, cream, terracotta, and olive support a premium textile house rather than a generic lifestyle label.", styles))
    brand_layout = Table(
        [[
            [
                Paragraph("Positioning", styles["CardTitle"]),
                Paragraph("Ali Dandin premium textiles", styles["BodyCopy"]),
                Paragraph("Voice", styles["CardTitle"]),
                Paragraph("Measured, exact, and materially intelligent. The brand speaks from judgment rather than hype.", styles["CardBody"]),
                Paragraph("System", styles["CardTitle"]),
                Paragraph("Primary mark, AD monogram, passport stamp, and textile globe anchor the visual architecture.", styles["CardBody"]),
            ],
            card([fitted_image(IMAGES / "logo_suite.png", 3.0 * inch, 3.6 * inch)], background=colors.HexColor("#F7F0E5"), padding=10),
        ]],
        colWidths=[3.45 * inch, 3.15 * inch],
    )
    brand_layout.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    palette_row = Table(
        [[
            stat_card("Midnight Blue", "#1E2A3A", styles),
            stat_card("Sand", "#C2A57A", styles),
            stat_card("Cream", "#FFFAF3", styles),
            stat_card("Terracotta / Olive", "Accent palette", styles),
        ]],
        colWidths=[1.65 * inch, 1.45 * inch, 1.45 * inch, 1.95 * inch],
    )
    palette_row.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(brand_layout)
    story.append(Spacer(1, 0.16 * inch))
    story.append(palette_row)
    story.append(PageBreak())

    story.extend(section_heading("Collection architecture", "Kyoto Indigo as the opening system", "Kyoto Indigo establishes the opening collection, combining material, dye process, and cultural origin into a coherent commercial system. The assortment is narrow by design so the brand remains legible and premium from day one.", styles))
    collection_layout = Table(
        [[
            card([fitted_image(IMAGES / "kyoto_indigo_collection.png", 2.7 * inch, 4.0 * inch)], background=colors.HexColor("#F7F0E5"), padding=10),
            [
                Paragraph("Collection logic", styles["CardTitle"]),
                Paragraph("One destination. One material story. One restrained assortment that can be photographed, merchandised, and repeated without dilution.", styles["CardBody"]),
                Spacer(1, 0.06 * inch),
                Paragraph("Opening assortment", styles["CardTitle"]),
            ]
        ]],
        colWidths=[2.9 * inch, 3.7 * inch],
    )
    collection_layout.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    product_grid = Table(
        [[
            narrative_card(product.name, product.note, styles, background=WHITE) for product in PRODUCTS[:2]
        ], [
            narrative_card(product.name, product.note, styles, background=WHITE) for product in PRODUCTS[2:]
        ]],
        colWidths=[3.18 * inch, 3.18 * inch],
    )
    product_grid.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(collection_layout)
    story.append(Spacer(1, 0.14 * inch))
    story.append(product_grid)
    story.append(PageBreak())

    story.extend(section_heading("Commerce platform", "Editorial storefront, disciplined hierarchy", "The storefront is collection-led rather than catalog-led. It presents one featured collection, one supporting story about process or maker, and one tightly edited product grid. That structure keeps the brand premium and intelligible.", styles))
    commerce_layout = Table(
        [[
            [
                Paragraph("Navigation", styles["CardTitle"]),
                Paragraph("Home, Collections, Journal, About, and Passport Club form the core public structure.", styles["CardBody"]),
                Paragraph("Merchandising", styles["CardTitle"]),
                Paragraph("The homepage leads with the collection hero, then reinforces the process story before moving into product cards.", styles["CardBody"]),
                Paragraph("Conversion", styles["CardTitle"]),
                Paragraph("Product pages elevate material, process, and origin so the brand sells judgment, not discounting.", styles["CardBody"]),
            ],
            card([fitted_image(IMAGES / "shopify_homepage.png", 2.9 * inch, 4.1 * inch)], background=colors.HexColor("#F7F0E5"), padding=10),
        ]],
        colWidths=[3.45 * inch, 3.15 * inch],
    )
    commerce_layout.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    commerce_cards = Table(
        [[
            narrative_card("Editorial engine", "Journal, short-form video, and launch stories create awareness without reducing the brand to performance marketing alone.", styles),
            narrative_card("Membership layer", "Passport Club creates recurring revenue and deepens early access, storytelling, and collection ownership.", styles),
            narrative_card("Private label path", "Once repeat demand is proven, signature textile lines support higher margin and broader distribution.", styles),
        ]],
        colWidths=[2.12 * inch, 2.12 * inch, 2.12 * inch],
    )
    commerce_cards.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(commerce_layout)
    story.append(Spacer(1, 0.14 * inch))
    story.append(commerce_cards)
    story.append(PageBreak())

    story.extend(section_heading("Sourcing model", "Regional depth and disciplined supplier logic", "Sourcing is a brand expression as much as an operating function. Japan establishes the opening story, while India, Turkey, Italy, and later Peru create a broader network across handcraft, scale, and luxury finishing.", styles))
    story.append(card([fitted_image(IMAGES / "global_sourcing_overview.png", 6.4 * inch, 3.6 * inch)], background=colors.HexColor("#F7F0E5"), padding=10))
    story.append(Spacer(1, 0.16 * inch))
    region_rows = []
    for idx in range(0, len(REGIONS), 2):
        pair = REGIONS[idx: idx + 2]
        region_rows.append([narrative_card(region.name, region.note, styles) for region in pair] + ([Spacer(0, 0)] if len(pair) == 1 else []))
    region_grid = Table(region_rows, colWidths=[3.18 * inch, 3.18 * inch])
    region_grid.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(region_grid)
    story.append(PageBreak())

    story.extend(section_heading("Financial architecture", "Premium pricing, disciplined assortment, layered revenue", "The base case grows from product-led revenue into a broader mix of membership, collaborations, and eventually private label. The model stays strongest when the assortment remains curated and the positioning remains textile-first.", styles))
    finance_layout = Table(
        [[financial_table(styles), unit_economics_table()]],
        colWidths=[4.0 * inch, 2.6 * inch],
    )
    finance_layout.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(finance_layout)
    story.append(Spacer(1, 0.18 * inch))
    finance_cards = Table(
        [[
            stat_card(money_short(GROSS_PROFIT[2]), "Year 3 gross profit", styles),
            stat_card(money_short(OPERATING_PROFIT[2]), "Year 3 operating profit", styles),
            stat_card("$8 / mo", "Membership pricing", styles),
        ]],
        colWidths=[2.1 * inch, 2.1 * inch, 2.1 * inch],
    )
    finance_cards.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(finance_cards)
    story.append(Spacer(1, 0.16 * inch))
    story.append(revenue_mix_table())
    story.append(PageBreak())

    story.extend(section_heading("Scale path", "From destination-led discovery into a modern textile house", "The company scales from discovery into authority and into systems that support long-term growth. The expansion path moves through repeat collections, private label, prestige partnerships, and eventual home textile breadth.", styles))
    story.append(roadmap_chart())
    story.append(Spacer(1, 0.18 * inch))
    scale_cards = Table(
        [[
            narrative_card("Phase 1", "Kyoto Indigo and subsequent destination-led collections establish the public point of view and commercial grammar.", styles),
            narrative_card("Phase 2", "Repeat collections, stronger supplier systems, and signature product lines create a brand with operating depth.", styles),
            narrative_card("Phase 3", "Bedding, home textiles, partnerships, licensing, and prestige retail extend the company into a lifestyle-scale brand.", styles),
        ]],
        colWidths=[2.12 * inch, 2.12 * inch, 2.12 * inch],
    )
    scale_cards.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(scale_cards)
    story.append(Spacer(1, 0.18 * inch))
    story.extend(section_heading("Capital strategy", "Tight early capital, stronger later leverage", "Initial capital supports product development, sourcing travel, content production, and ecommerce operations. Later capital supports inventory scale, team build-out, and broader supplier management after the brand has already established proof.", styles))
    capital_cards = Table(
        [[
            narrative_card("Use of funds", "Product development, sourcing, inventory, content, and ecommerce operations.", styles),
            narrative_card("Capital range", "$35K-$70K for the initial operating phase, followed by growth capital after traction.", styles),
            narrative_card("Long-term upside", "A premium textile authority with room to scale into signature home collections, licensing, and selective retail.", styles),
        ]],
        colWidths=[2.12 * inch, 2.12 * inch, 2.12 * inch],
    )
    capital_cards.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(capital_cards)
    story.append(Spacer(1, 0.18 * inch))
    story.append(capital_plan_table())
    story.append(PageBreak())

    story.extend(section_heading("Closing perspective", "A premium textile brand with room to compound", "Ali Dandin is strongest when it stays narrow, exacting, and textile-led. The brand earns authority through destination-based collections, premium product judgment, and a commercial system that compounds through repeat launches and signature lines.", styles))
    close_layout = Table(
        [[
            [
                Paragraph("What lasts", styles["CardTitle"]),
                Paragraph("The long-term asset is customer trust around taste, provenance, and material literacy.", styles["CardBody"]),
                Spacer(1, 0.08 * inch),
                Paragraph("What scales", styles["CardTitle"]),
                Paragraph("Collections, membership, private label, and prestige partnerships expand the platform without breaking the core brand thesis.", styles["CardBody"]),
                Spacer(1, 0.08 * inch),
                Paragraph("What differentiates", styles["CardTitle"]),
                Paragraph("Few brands combine textile fluency, destination-led curation, and premium editorial commerce in one coherent system.", styles["CardBody"]),
            ],
            card([fitted_image(IMAGES / "shopify_homepage.png", 3.05 * inch, 4.25 * inch)], background=colors.HexColor("#F7F0E5"), padding=10),
        ]],
        colWidths=[3.55 * inch, 3.05 * inch],
    )
    close_layout.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    close_stats = Table(
        [[
            stat_card(money_short(PRODUCT_REVENUE[-1]), "Year 3 product revenue", styles),
            stat_card(f"{int(GROSS_MARGIN[-1] * 100)}%", "Indicative gross margin", styles),
            stat_card(f"{MEMBERSHIP_COUNT[-1]:,}", "Year 3 members", styles),
        ]],
        colWidths=[2.1 * inch, 2.1 * inch, 2.1 * inch],
    )
    close_stats.setStyle(TableStyle([("LEFTPADDING", (0, 0), (-1, -1), 0), ("RIGHTPADDING", (0, 0), (-1, -1), 0)]))
    story.append(close_layout)
    story.append(Spacer(1, 0.16 * inch))
    story.append(close_stats)

    doc = SimpleDocTemplate(
        str(output),
        pagesize=letter,
        rightMargin=0.58 * inch,
        leftMargin=0.58 * inch,
        topMargin=0.55 * inch,
        bottomMargin=0.55 * inch,
        title="Ali Dandin Business Plan",
        author="OpenAI Codex",
    )
    doc.build(story, onFirstPage=draw_first_page, onLaterPages=draw_later_pages)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size, color, *, bold=False, font_name="Aptos", align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    frame = tx.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tx


def add_panel(slide, left, top, width, height, *, fill_rgb, line_rgb, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    return shape


def add_picture_contain(slide, path: Path, left, top, width, height):
    image_w, image_h = image_dimensions(path)
    frame_ratio = width / height
    image_ratio = image_w / image_h
    if image_ratio > frame_ratio:
        draw_width = width
        draw_height = width / image_ratio
        draw_left = left
        draw_top = top + (height - draw_height) / 2
    else:
        draw_height = height
        draw_width = height * image_ratio
        draw_left = left + (width - draw_width) / 2
        draw_top = top
    slide.shapes.add_picture(str(path), draw_left, draw_top, width=draw_width, height=draw_height)


def add_bullet_card(slide, left, top, width, height, title, bullets, *, fill_rgb, line_rgb, title_rgb, body_rgb):
    panel = add_panel(slide, left, top, width, height, fill_rgb=fill_rgb, line_rgb=line_rgb)
    frame = panel.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = title_rgb
    p.font.name = "Georgia"
    for bullet in bullets:
        para = frame.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(11.5)
        para.font.color.rgb = body_rgb
        para.font.name = "Aptos"
        para.space_after = Pt(5)
        para.bullet = True
    return panel


def add_stat_box(slide, left, top, width, height, value, label, *, fill_rgb, line_rgb, value_rgb, label_rgb):
    panel = add_panel(slide, left, top, width, height, fill_rgb=fill_rgb, line_rgb=line_rgb)
    frame = panel.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = value_rgb
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER
    p = frame.add_paragraph()
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = label_rgb
    p.font.name = "Aptos"
    p.alignment = PP_ALIGN.CENTER
    return panel


def add_note_card(slide, left, top, width, height, kicker, title, body, *, fill_rgb, line_rgb, kicker_rgb, title_rgb, body_rgb):
    panel = add_panel(slide, left, top, width, height, fill_rgb=fill_rgb, line_rgb=line_rgb)
    frame = panel.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = kicker
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = kicker_rgb
    p.font.name = "Aptos"
    p = frame.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = title_rgb
    p.font.name = "Georgia"
    p.space_before = Pt(4)
    p = frame.add_paragraph()
    p.text = body
    p.font.size = Pt(10.5)
    p.font.color.rgb = body_rgb
    p.font.name = "Aptos"
    p.space_before = Pt(5)
    return panel


def add_revenue_chart_slide(slide, left, top, width, height, *, bg_rgb, axis_rgb, title_rgb):
    add_panel(slide, left, top, width, height, fill_rgb=bg_rgb, line_rgb=RGBColor(0xD8, 0xCF, 0xC2))
    add_text(slide, left + Inches(0.18), top + Inches(0.08), width - Inches(0.36), Inches(0.25), "Revenue ramp", 13, title_rgb, bold=True, font_name="Aptos")
    max_value = max(REVENUE)
    chart_top = top + Inches(0.48)
    chart_bottom = top + height - Inches(0.5)
    chart_height = chart_bottom - chart_top
    bar_width = Inches(0.68)
    gap = Inches(0.56)
    base_left = left + Inches(0.55)
    line_y = chart_bottom
    axis_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.42), line_y, width - Inches(0.84), Inches(0.02))
    axis_line.fill.solid()
    axis_line.fill.fore_color.rgb = axis_rgb
    axis_line.line.color.rgb = axis_rgb
    for idx, (label, value) in enumerate(zip(YEAR_LABELS, REVENUE)):
        bar_height = chart_height * (value / max_value)
        x = base_left + idx * (bar_width + gap)
        y = chart_bottom - bar_height
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, bar_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY_RGB
        bar.line.color.rgb = NAVY_RGB
        add_text(slide, x - Inches(0.05), y - Inches(0.28), bar_width + Inches(0.1), Inches(0.2), money_short(value), 10, title_rgb, bold=True, font_name="Aptos", align=PP_ALIGN.CENTER)
        add_text(slide, x - Inches(0.1), chart_bottom + Inches(0.05), bar_width + Inches(0.2), Inches(0.22), label, 10, axis_rgb, font_name="Aptos", align=PP_ALIGN.CENTER)


NAVY_RGB = RGBColor(0x1E, 0x2A, 0x3A)
SAND_RGB = RGBColor(0xC2, 0xA5, 0x7A)
CREAM_RGB = RGBColor(0xFF, 0xFA, 0xF3)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_RGB = RGBColor(0xF7, 0xF0, 0xE5)
LINE_RGB = RGBColor(0xD8, 0xCF, 0xC2)
MUTED_RGB = RGBColor(0x65, 0x70, 0x7A)
TERRACOTTA_RGB = RGBColor(0xB9, 0x69, 0x39)
OLIVE_RGB = RGBColor(0x61, 0x71, 0x47)


def add_slide_header(slide, page_number: int, title: str, subtitle: str, *, dark: bool = False):
    kicker_color = SAND_RGB if dark else TERRACOTTA_RGB
    title_color = CREAM_RGB if dark else NAVY_RGB
    subtitle_color = CREAM_RGB if dark else MUTED_RGB
    add_text(slide, Inches(0.7), Inches(0.42), Inches(1.0), Inches(0.2), f"{page_number:02d}", 10, kicker_color, bold=True)
    add_text(slide, Inches(0.7), Inches(0.78), Inches(5.9), Inches(0.55), title, 27 if dark else 25, title_color, bold=True, font_name="Georgia")
    add_text(slide, Inches(0.7), Inches(1.48), Inches(5.95), Inches(0.72), subtitle, 13.5, subtitle_color, font_name="Aptos")


def build_pptx() -> None:
    output = DOCS / "investor_deck.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY_RGB)
    add_panel(slide, Inches(7.72), Inches(0.46), Inches(4.88), Inches(6.56), fill_rgb=RGBColor(0x2A, 0x39, 0x4C), line_rgb=RGBColor(0x2A, 0x39, 0x4C))
    add_picture_contain(slide, IMAGES / "kyoto_indigo_collection.png", Inches(7.86), Inches(0.58), Inches(4.6), Inches(6.32))
    add_text(slide, Inches(0.75), Inches(0.64), Inches(2.4), Inches(0.22), "ALI DANDIN", 11, SAND_RGB, bold=True)
    add_text(slide, Inches(0.75), Inches(1.15), Inches(5.7), Inches(1.2), "Ali Dandin premium textiles", 31, CREAM_RGB, bold=True, font_name="Georgia")
    add_text(slide, Inches(0.75), Inches(2.3), Inches(5.7), Inches(1.0), "A destination-led textile brand built around provenance, material intelligence, and premium direct-to-consumer commerce.", 15, CREAM_RGB, font_name="Aptos")
    add_stat_box(slide, Inches(0.75), Inches(4.45), Inches(1.35), Inches(1.02), "$826K", "Year 1 revenue", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x3F, 0x52, 0x6B), value_rgb=CREAM_RGB, label_rgb=RGBColor(0xDA, 0xD4, 0xCA))
    add_stat_box(slide, Inches(2.25), Inches(4.45), Inches(1.35), Inches(1.02), "4", "Opening SKUs", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x3F, 0x52, 0x6B), value_rgb=CREAM_RGB, label_rgb=RGBColor(0xDA, 0xD4, 0xCA))
    add_stat_box(slide, Inches(3.75), Inches(4.45), Inches(1.35), Inches(1.02), "55-65%", "Gross margin", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x3F, 0x52, 0x6B), value_rgb=CREAM_RGB, label_rgb=RGBColor(0xDA, 0xD4, 0xCA))
    add_stat_box(slide, Inches(5.25), Inches(4.45), Inches(1.35), Inches(1.02), "Kyoto", "Launch story", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x3F, 0x52, 0x6B), value_rgb=CREAM_RGB, label_rgb=RGBColor(0xDA, 0xD4, 0xCA))
    add_text(slide, Inches(0.75), Inches(6.18), Inches(5.9), Inches(0.55), "A category with pricing power, repeatability, and room to expand into signature lines.", 12.5, RGBColor(0xD8, 0xD1, 0xC7), font_name="Aptos")

    # Market gap
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM_RGB)
    add_slide_header(slide, 2, "Market gap", "Premium buyers want provenance, but current channels flatten textiles into generic product.")
    add_note_card(slide, Inches(0.72), Inches(2.18), Inches(3.78), Inches(4.42), "01", "Too much catalog", "Selection breadth overwhelms curation, so material judgment is buried inside generic assortment.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(4.78), Inches(2.18), Inches(3.78), Inches(4.42), "02", "Premium without depth", "High-end positioning exists, but often without sourcing credibility or a clear textile-specific thesis.", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, kicker_rgb=TERRACOTTA_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(8.84), Inches(2.18), Inches(3.78), Inches(4.42), "03", "Fragmented trust", "Customers still lack a single textile authority they rely on for quality, taste, and origin.", fill_rgb=RGBColor(0x1E, 0x2A, 0x3A), line_rgb=RGBColor(0x32, 0x43, 0x59), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))

    # Brand position
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM_RGB)
    add_slide_header(slide, 3, "Brand position", "Ali Dandin is a premium textile brand built around provenance, material intelligence, and destination-led collections.")
    add_panel(slide, Inches(0.72), Inches(2.1), Inches(5.6), Inches(4.7), fill_rgb=WHITE_RGB, line_rgb=LINE_RGB)
    add_text(slide, Inches(0.98), Inches(2.38), Inches(4.7), Inches(0.35), "Brand position", 19, NAVY_RGB, bold=True, font_name="Georgia")
    add_text(slide, Inches(0.98), Inches(2.82), Inches(4.9), Inches(0.8), "The offer narrows to what is ownable: textiles, place, process, and curation.", 13, MUTED_RGB)
    add_stat_box(slide, Inches(0.98), Inches(4.0), Inches(1.25), Inches(1.02), "1", "Category", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(2.38), Inches(4.0), Inches(1.25), Inches(1.02), "1", "Collection", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(3.78), Inches(4.0), Inches(1.25), Inches(1.02), "4", "Core marks", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_text(slide, Inches(0.98), Inches(5.32), Inches(1.6), Inches(0.22), "Voice", 11.5, TERRACOTTA_RGB, bold=True)
    add_text(slide, Inches(0.98), Inches(5.58), Inches(4.9), Inches(0.32), "Measured, materially intelligent, and anchored in provenance.", 11.2, NAVY_RGB)
    add_panel(slide, Inches(6.58), Inches(2.1), Inches(6.02), Inches(4.7), fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB)
    add_picture_contain(slide, IMAGES / "logo_suite.png", Inches(6.72), Inches(2.24), Inches(5.74), Inches(4.42))

    # Opening collection
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY_RGB)
    add_slide_header(slide, 4, "Kyoto Indigo", "The opening collection defines how the brand will look, merchandise, and scale.", dark=True)
    add_panel(slide, Inches(0.72), Inches(2.08), Inches(4.15), Inches(4.76), fill_rgb=RGBColor(0x27, 0x36, 0x49), line_rgb=RGBColor(0x35, 0x47, 0x5E))
    add_picture_contain(slide, IMAGES / "kyoto_indigo_collection.png", Inches(0.86), Inches(2.22), Inches(3.87), Inches(4.48))
    add_bullet_card(slide, Inches(5.1), Inches(2.08), Inches(3.45), Inches(2.22), "Opening assortment", [
        "Indigo Linen Throw",
        "Indigo Scarf",
        "Indigo Pillow",
        "Indigo Textile Set",
    ], fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_bullet_card(slide, Inches(8.8), Inches(2.08), Inches(3.8), Inches(2.22), "Collection logic", [
        "One destination, one material story, one disciplined assortment.",
        "Clear photography, merchandising, and memory structure.",
        "Strong foundation for future destination-led releases.",
    ], fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_stat_box(slide, Inches(5.1), Inches(4.62), Inches(2.28), Inches(1.12), "4", "Opening SKUs", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(7.62), Inches(4.62), Inches(2.28), Inches(1.12), "Japan", "Origin story", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(10.14), Inches(4.62), Inches(2.42), Inches(1.12), "Textile-led", "Category position", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)

    # Commerce platform
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM_RGB)
    add_slide_header(slide, 5, "Commerce platform", "The storefront is editorial, collection-led, and designed to convert on taste and provenance.")
    add_panel(slide, Inches(0.72), Inches(2.08), Inches(4.4), Inches(4.86), fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB)
    add_picture_contain(slide, IMAGES / "shopify_homepage.png", Inches(0.86), Inches(2.22), Inches(4.12), Inches(4.58))
    add_note_card(slide, Inches(5.38), Inches(2.08), Inches(3.4), Inches(2.25), "01", "Collection hero", "The homepage opens with one featured collection instead of a broad catalog grid.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(8.98), Inches(2.08), Inches(3.58), Inches(2.25), "02", "Process story", "Editorial blocks explain the material, maker, and origin before hard selling the assortment.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(5.38), Inches(4.69), Inches(3.4), Inches(2.25), "03", "Product clarity", "Collection and product pages are structured around material, process, and provenance.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(8.98), Inches(4.69), Inches(3.58), Inches(2.25), "04", "Membership layer", "Passport Club sits inside the brand as a deeper relationship rather than a disruptive conversion trick.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)

    # Sourcing network
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY_RGB)
    add_slide_header(slide, 6, "Sourcing network", "Regional sourcing supports both product depth and brand meaning.", dark=True)
    add_panel(slide, Inches(7.02), Inches(1.95), Inches(5.56), Inches(4.9), fill_rgb=RGBColor(0x27, 0x36, 0x49), line_rgb=RGBColor(0x35, 0x47, 0x5E))
    add_picture_contain(slide, IMAGES / "global_sourcing_overview.png", Inches(7.16), Inches(2.08), Inches(5.28), Inches(4.62))
    add_bullet_card(slide, Inches(0.72), Inches(1.95), Inches(5.96), Inches(1.72), "Regional priorities", [
        "Japan: heritage process and indigo authority",
        "India: handloom and block print expansion",
        "Turkey and Italy: scale and premium finishing",
    ], fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_stat_box(slide, Inches(0.72), Inches(4.04), Inches(1.8), Inches(1.05), "Japan", "Process anchor", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(2.66), Inches(4.04), Inches(1.8), Inches(1.05), "India", "Breadth", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(4.6), Inches(4.04), Inches(1.8), Inches(1.05), "Turkey", "Scale", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(0.72), Inches(5.36), Inches(1.8), Inches(1.05), "Italy", "Finishing", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(2.66), Inches(5.36), Inches(1.8), Inches(1.05), "Peru", "Future fiber", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_text(slide, Inches(4.88), Inches(5.43), Inches(1.9), Inches(0.82), "Supplier depth follows the same rule as the assortment: narrow first, broader second.", 12, RGBColor(0xD8, 0xD1, 0xC7), font_name="Aptos")

    # Revenue stack
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM_RGB)
    add_slide_header(slide, 7, "Revenue stack", "Product sales lead; membership and partnerships deepen the model before broader private label expansion.")
    add_note_card(slide, Inches(0.72), Inches(2.08), Inches(3.78), Inches(2.24), "01", "Collection sales", "Destination-led drops remain the core engine, representing the majority of Year 3 revenue and reinforcing the textile-first thesis.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(4.78), Inches(2.08), Inches(3.78), Inches(2.24), "02", "Membership", "Passport Club adds recurring revenue, improves retention, and deepens ownership of launches through early access and editorial benefits.", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, kicker_rgb=TERRACOTTA_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(8.84), Inches(2.08), Inches(3.78), Inches(2.24), "03", "Partnerships", "Selective collaborations and affiliate income extend monetization without turning the brand into a broad wholesale platform.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(0.72), Inches(4.76), Inches(2.38), Inches(1.18), "$9.35M", "Year 3 total revenue", fill_rgb=RGBColor(0x1E, 0x2A, 0x3A), line_rgb=RGBColor(0x1E, 0x2A, 0x3A), value_rgb=CREAM_RGB, label_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_stat_box(slide, Inches(3.38), Inches(4.76), Inches(2.38), Inches(1.18), "74%", "Product mix", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(6.04), Inches(4.76), Inches(2.38), Inches(1.18), "15%", "Membership mix", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(8.7), Inches(4.76), Inches(2.38), Inches(1.18), "11%", "Partnership mix", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(11.36), Inches(4.76), Inches(1.26), Inches(1.18), "PL", "Future layer", fill_rgb=RGBColor(0x61, 0x71, 0x47), line_rgb=RGBColor(0x61, 0x71, 0x47), value_rgb=CREAM_RGB, label_rgb=RGBColor(0xF0, 0xEA, 0xDF))

    # Economic model
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY_RGB)
    add_slide_header(slide, 8, "Economic model", "Premium ASPs and disciplined assortment produce a clean revenue ramp.", dark=True)
    add_revenue_chart_slide(slide, Inches(0.72), Inches(2.02), Inches(5.76), Inches(4.9), bg_rgb=CREAM_RGB, axis_rgb=MUTED_RGB, title_rgb=NAVY_RGB)
    add_stat_box(slide, Inches(6.8), Inches(2.12), Inches(2.6), Inches(1.12), "$6.91M", "Year 3 product revenue", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(9.68), Inches(2.12), Inches(2.88), Inches(1.12), "57-61%", "Gross margin range", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(6.8), Inches(3.55), Inches(2.6), Inches(1.12), "12", "Collections by Year 3", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(9.68), Inches(3.55), Inches(2.88), Inches(1.12), "$8", "Membership price", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_note_card(slide, Inches(6.8), Inches(5.0), Inches(5.76), Inches(1.9), "WHY IT HOLDS", "Clean revenue ramp", "Small-batch launches reduce inventory exposure while preserving premium brand specificity and room for private label expansion.", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))

    # Scale path
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM_RGB)
    add_slide_header(slide, 9, "Scale path", "The company advances from discovery into authority and then into a broader textile platform.")
    add_panel(slide, Inches(0.72), Inches(2.08), Inches(11.86), Inches(2.28), fill_rgb=WHITE_RGB, line_rgb=LINE_RGB)
    line_y = Inches(3.1)
    timeline_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.35), line_y, Inches(10.56), Inches(0.03))
    timeline_line.fill.solid()
    timeline_line.fill.fore_color.rgb = SAND_RGB
    timeline_line.line.color.rgb = SAND_RGB
    nodes = [
        (Inches(2.1), "Discovery", "Years 1-2", "Destination-led launches"),
        (Inches(6.05), "Authority", "Years 3-4", "Repeat collections and private label"),
        (Inches(10.0), "Expansion", "Years 5+", "Home textiles, partnerships, licensing"),
    ]
    for x, title, period, note in nodes:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, line_y - Inches(0.17), Inches(0.34), Inches(0.34))
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY_RGB
        circ.line.color.rgb = NAVY_RGB
        add_text(slide, x - Inches(0.28), Inches(2.38), Inches(1.05), Inches(0.22), title, 12, NAVY_RGB, bold=True, font_name="Aptos", align=PP_ALIGN.CENTER)
        add_text(slide, x - Inches(0.28), Inches(2.62), Inches(1.05), Inches(0.2), period, 10, TERRACOTTA_RGB, font_name="Aptos", align=PP_ALIGN.CENTER)
        add_text(slide, x - Inches(0.7), Inches(3.45), Inches(1.9), Inches(0.45), note, 10.5, MUTED_RGB, font_name="Aptos", align=PP_ALIGN.CENTER)
    add_bullet_card(slide, Inches(0.72), Inches(4.7), Inches(3.74), Inches(2.0), "Phase 1", [
        "Clear imagery",
        "Strong assortment discipline",
        "Focused destination story",
    ], fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_bullet_card(slide, Inches(4.78), Inches(4.7), Inches(3.74), Inches(2.0), "Phase 2", [
        "Supplier depth",
        "Signature lines",
        "Private label confidence",
    ], fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_bullet_card(slide, Inches(8.84), Inches(4.7), Inches(3.74), Inches(2.0), "Phase 3", [
        "Prestige retail",
        "Licensing",
        "Lifestyle-scale brand equity",
    ], fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)

    # Capital use / close
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY_RGB)
    add_slide_header(slide, 10, "Capital use", "Early capital funds product, sourcing, content, and operating depth; later capital supports scale.", dark=True)
    add_note_card(slide, Inches(0.72), Inches(2.08), Inches(2.8), Inches(2.18), "01", "Product", "Inventory, sampling, and development for a narrow premium assortment.", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_note_card(slide, Inches(3.78), Inches(2.08), Inches(2.8), Inches(2.18), "02", "Sourcing", "Travel, supplier onboarding, and quality control across priority regions.", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_note_card(slide, Inches(0.72), Inches(4.62), Inches(2.8), Inches(2.18), "03", "Content", "Editorial production that turns process and origin into conversion-ready story.", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_note_card(slide, Inches(3.78), Inches(4.62), Inches(2.8), Inches(2.18), "04", "Operations", "Ecommerce, fulfillment, and early team support to keep launches reliable.", fill_rgb=RGBColor(0x2B, 0x3A, 0x4C), line_rgb=RGBColor(0x35, 0x47, 0x5E), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_stat_box(slide, Inches(6.94), Inches(2.1), Inches(5.48), Inches(1.35), "$35K-$70K", "Initial capital range", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_note_card(slide, Inches(6.94), Inches(3.86), Inches(5.48), Inches(2.94), "INVESTOR CASE", "Why this works", "Ali Dandin combines category clarity, sourcing fluency, premium pricing power, and a disciplined path from edited collections into signature lines and selective partnerships.", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=RGBColor(0x35, 0x47, 0x5E), kicker_rgb=TERRACOTTA_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_text(slide, Inches(0.72), Inches(6.92), Inches(11.86), Inches(0.2), "Ali Dandin premium textiles", 15, SAND_RGB, bold=True, font_name="Georgia", align=PP_ALIGN.CENTER)

    # Close
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM_RGB)
    add_text(slide, Inches(0.88), Inches(0.72), Inches(2.4), Inches(0.22), "ALI DANDIN", 11, TERRACOTTA_RGB, bold=True)
    add_text(slide, Inches(0.88), Inches(1.26), Inches(7.0), Inches(1.0), "Ali Dandin premium textiles", 32, NAVY_RGB, bold=True, font_name="Georgia")
    add_text(slide, Inches(0.88), Inches(2.24), Inches(7.2), Inches(1.1), "A premium textile brand built around provenance, material intelligence, and destination-led collections.", 16, MUTED_RGB, font_name="Aptos")
    add_note_card(slide, Inches(0.88), Inches(3.58), Inches(3.55), Inches(2.08), "WHY IT MATTERS", "Category trust", "The brand owns a clear point of view in textiles rather than competing as a broad lifestyle or travel catalog.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(4.62), Inches(3.58), Inches(3.55), Inches(2.08), "WHY IT SCALES", "Repeatable system", "Collections, membership, private label, and partnerships extend the same thesis without breaking it.", fill_rgb=WHITE_RGB, line_rgb=LINE_RGB, kicker_rgb=SAND_RGB, title_rgb=NAVY_RGB, body_rgb=MUTED_RGB)
    add_note_card(slide, Inches(8.36), Inches(3.58), Inches(3.55), Inches(2.08), "WHY IT WINS", "Material authority", "Ali's textile fluency turns product judgment and sourcing discipline into pricing power and brand trust.", fill_rgb=RGBColor(0x1E, 0x2A, 0x3A), line_rgb=RGBColor(0x1E, 0x2A, 0x3A), kicker_rgb=SAND_RGB, title_rgb=CREAM_RGB, body_rgb=RGBColor(0xDD, 0xD6, 0xCC))
    add_stat_box(slide, Inches(0.88), Inches(6.1), Inches(2.25), Inches(0.94), "$9.35M", "Year 3 revenue", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(3.38), Inches(6.1), Inches(2.25), Inches(0.94), "55-65%", "Gross margin", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_stat_box(slide, Inches(5.88), Inches(6.1), Inches(2.25), Inches(0.94), "12", "Collections by Year 3", fill_rgb=RGBColor(0xF7, 0xF0, 0xE5), line_rgb=LINE_RGB, value_rgb=NAVY_RGB, label_rgb=MUTED_RGB)
    add_text(slide, Inches(8.56), Inches(6.28), Inches(3.4), Inches(0.28), "Premium textiles shaped by place, process, and trust.", 13, NAVY_RGB, bold=True, font_name="Georgia")

    prs.core_properties.title = "Ali Dandin Investor Deck"
    prs.core_properties.author = "OpenAI Codex"
    prs.save(str(output))


def build_xlsx() -> None:
    output = DOCS / "financial_model.xlsx"
    wb = Workbook()
    wb.remove(wb.active)

    navy_fill = PatternFill("solid", fgColor="1E2A3A")
    sand_fill = PatternFill("solid", fgColor="F2E7D6")
    green_fill = PatternFill("solid", fgColor="EEF4E7")
    border = Border(
        left=Side(style="thin", color="D7D0C5"),
        right=Side(style="thin", color="D7D0C5"),
        top=Side(style="thin", color="D7D0C5"),
        bottom=Side(style="thin", color="D7D0C5"),
    )

    def style_header(cell, dark=True):
        cell.font = Font(bold=True, color="FFFFFF" if dark else "1E2A3A")
        cell.fill = navy_fill if dark else sand_fill
        cell.border = border
        cell.alignment = Alignment(horizontal="center", vertical="center")

    def style_cell(cell, bold=False, fill=None, number_format=None):
        cell.font = Font(bold=bold, color="1F2A36")
        if fill:
            cell.fill = fill
        cell.border = border
        if number_format:
            cell.number_format = number_format
        cell.alignment = Alignment(vertical="center")

    def set_widths(ws, widths):
        for idx, width in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(idx)].width = width

    ws = wb.create_sheet("Assumptions")
    ws.append(["Metric", "Year 1", "Year 2", "Year 3", "Notes"])
    assumption_rows = [
        ("Audience size", AUDIENCE[0], AUDIENCE[1], AUDIENCE[2], "Combined social, email, and owned audience"),
        ("Collections launched", COLLECTIONS[0], COLLECTIONS[1], COLLECTIONS[2], "Quarterly to monthly cadence"),
        ("Total units sold", TOTAL_UNITS[0], TOTAL_UNITS[1], TOTAL_UNITS[2], "Scaled through launch cadence and repeat purchase"),
        ("Average selling price", AVERAGE_SELLING_PRICE[0], AVERAGE_SELLING_PRICE[1], AVERAGE_SELLING_PRICE[2], "Premium textile pricing"),
        ("Product revenue", PRODUCT_REVENUE[0], PRODUCT_REVENUE[1], PRODUCT_REVENUE[2], "Collection and product sales"),
        ("Membership conversion", MEMBERSHIP_CONVERSION[0], MEMBERSHIP_CONVERSION[1], MEMBERSHIP_CONVERSION[2], "Audience conversion into Passport Club"),
        ("Membership count", MEMBERSHIP_COUNT[0], MEMBERSHIP_COUNT[1], MEMBERSHIP_COUNT[2], "Recurring member base"),
        ("Membership price", MEMBERSHIP_PRICE, MEMBERSHIP_PRICE, MEMBERSHIP_PRICE, "Monthly Passport Club price"),
        ("Membership revenue", MEMBERSHIP_REVENUE[0], MEMBERSHIP_REVENUE[1], MEMBERSHIP_REVENUE[2], "Annual recurring revenue"),
        ("Partnership revenue", PARTNERSHIP_REVENUE[0], PARTNERSHIP_REVENUE[1], PARTNERSHIP_REVENUE[2], "Affiliate and brand partnership income"),
        ("Total revenue", REVENUE[0], REVENUE[1], REVENUE[2], "Combined revenue view"),
        ("COGS", COGS[0], COGS[1], COGS[2], "Product, import, packaging, and fulfillment"),
        ("Travel & sourcing", TRAVEL_AND_SOURCING[0], TRAVEL_AND_SOURCING[1], TRAVEL_AND_SOURCING[2], "Trips, supplier meetings, and quality control"),
        ("Team cost", TEAM_COST[0], TEAM_COST[1], TEAM_COST[2], "Core staffing and contractor support"),
        ("Marketing & growth", MARKETING_AND_GROWTH[0], MARKETING_AND_GROWTH[1], MARKETING_AND_GROWTH[2], "Content amplification and demand generation"),
        ("General operations", GENERAL_OPS[0], GENERAL_OPS[1], GENERAL_OPS[2], "Tools, admin, platform, and overhead"),
        ("Total operating costs", OPERATING_COSTS[0], OPERATING_COSTS[1], OPERATING_COSTS[2], "Operating expense view"),
        ("Gross profit", GROSS_PROFIT[0], GROSS_PROFIT[1], GROSS_PROFIT[2], "Revenue minus COGS"),
        ("Operating profit", OPERATING_PROFIT[0], OPERATING_PROFIT[1], OPERATING_PROFIT[2], "Gross profit minus operating costs"),
    ]
    for row in assumption_rows:
        ws.append(list(row))
    for cell in ws[1]:
        style_header(cell)
    currency_rows = {
        "Average selling price",
        "Product revenue",
        "Membership price",
        "Membership revenue",
        "Partnership revenue",
        "Total revenue",
        "COGS",
        "Travel & sourcing",
        "Team cost",
        "Marketing & growth",
        "General operations",
        "Total operating costs",
        "Gross profit",
        "Operating profit",
    }
    percent_rows = {"Membership conversion"}
    highlighted_rows = {"Total revenue", "Gross profit", "Operating profit"}
    for row in ws.iter_rows(min_row=2):
        label = row[0].value
        for idx, cell in enumerate(row):
            style_cell(cell, bold=label in highlighted_rows, fill=green_fill if label in highlighted_rows else None)
            if 1 <= idx <= 3:
                if label in percent_rows:
                    cell.number_format = "0.0%"
                elif label in currency_rows:
                    cell.number_format = "$#,##0"
                else:
                    cell.number_format = "#,##0"
    set_widths(ws, [24, 15, 15, 15, 44])

    ws = wb.create_sheet("SKU_Model")
    sku_header = [
        "SKU",
        "Unit Mix %",
        "Role",
        "Y1 Units",
        "Y1 ASP",
        "Y1 Revenue",
        "Y1 Cost/Unit",
        "Y1 Gross Profit",
        "Y2 Units",
        "Y2 ASP",
        "Y2 Revenue",
        "Y2 Cost/Unit",
        "Y2 Gross Profit",
        "Y3 Units",
        "Y3 ASP",
        "Y3 Revenue",
        "Y3 Cost/Unit",
        "Y3 Gross Profit",
    ]
    ws.append(sku_header)
    for row_idx, sku in enumerate(SKU_MODELS, start=2):
        ws.cell(row=row_idx, column=1, value=sku.name)
        ws.cell(row=row_idx, column=2, value=sku.unit_mix)
        ws.cell(row=row_idx, column=3, value=sku.role)
        ws.cell(row=row_idx, column=4, value=f"=ROUND(Assumptions!B4*$B{row_idx},0)")
        ws.cell(row=row_idx, column=5, value=sku.asp[0])
        ws.cell(row=row_idx, column=6, value=f"=D{row_idx}*E{row_idx}")
        ws.cell(row=row_idx, column=7, value=sku.direct_cost[0])
        ws.cell(row=row_idx, column=8, value=f"=F{row_idx}-(D{row_idx}*G{row_idx})")
        ws.cell(row=row_idx, column=9, value=f"=ROUND(Assumptions!C4*$B{row_idx},0)")
        ws.cell(row=row_idx, column=10, value=sku.asp[1])
        ws.cell(row=row_idx, column=11, value=f"=I{row_idx}*J{row_idx}")
        ws.cell(row=row_idx, column=12, value=sku.direct_cost[1])
        ws.cell(row=row_idx, column=13, value=f"=K{row_idx}-(I{row_idx}*L{row_idx})")
        ws.cell(row=row_idx, column=14, value=f"=ROUND(Assumptions!D4*$B{row_idx},0)")
        ws.cell(row=row_idx, column=15, value=sku.asp[2])
        ws.cell(row=row_idx, column=16, value=f"=N{row_idx}*O{row_idx}")
        ws.cell(row=row_idx, column=17, value=sku.direct_cost[2])
        ws.cell(row=row_idx, column=18, value=f"=P{row_idx}-(N{row_idx}*Q{row_idx})")
    total_row = len(SKU_MODELS) + 2
    ws.cell(row=total_row, column=1, value="Illustrative total")
    for col in [4, 6, 8, 9, 11, 13, 14, 16, 18]:
        letter = get_column_letter(col)
        ws.cell(row=total_row, column=col, value=f"=SUM({letter}2:{letter}{total_row - 1})")
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=total_row):
        is_total = row[0].row == total_row
        for idx, cell in enumerate(row, start=1):
            style_cell(cell, bold=is_total, fill=green_fill if is_total else None)
            if idx == 2:
                cell.number_format = "0.0%"
            elif idx in {5, 6, 7, 8, 10, 11, 12, 13, 15, 16, 17, 18}:
                cell.number_format = "$#,##0"
            elif idx in {4, 9, 14}:
                cell.number_format = "#,##0"
    set_widths(ws, [24, 11, 34, 12, 10, 13, 11, 13, 12, 10, 13, 11, 13, 12, 10, 13, 11, 13])

    ws = wb.create_sheet("Launch_Cadence")
    ws.append(["Year", "Launch count", "Featured collections", "Avg units / launch", "Avg revenue / launch", "Operating focus"])
    for row_idx, (year_label, cadence, featured, note) in enumerate(LAUNCH_CADENCE, start=2):
        year_number = row_idx - 1
        ws.cell(row=row_idx, column=1, value=year_label)
        ws.cell(row=row_idx, column=2, value=f"=Assumptions!{get_column_letter(year_number + 1)}3")
        ws.cell(row=row_idx, column=3, value=featured)
        ws.cell(row=row_idx, column=4, value=f"=Assumptions!{get_column_letter(year_number + 1)}4/Assumptions!{get_column_letter(year_number + 1)}3")
        ws.cell(row=row_idx, column=5, value=f"=Assumptions!{get_column_letter(year_number + 1)}6/Assumptions!{get_column_letter(year_number + 1)}3")
        ws.cell(row=row_idx, column=6, value=note)
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=4):
        for idx, cell in enumerate(row, start=1):
            style_cell(cell)
            if idx == 4:
                cell.number_format = "#,##0"
            elif idx == 5:
                cell.number_format = "$#,##0"
    set_widths(ws, [12, 12, 42, 16, 18, 50])

    ws = wb.create_sheet("Membership_Model")
    ws.append(["Metric", "Year 1", "Year 2", "Year 3", "Method"])
    membership_rows = [
        ("Audience size", "=Assumptions!B2", "=Assumptions!C2", "=Assumptions!D2", "Pulled from assumptions"),
        ("Conversion rate", "=Assumptions!B7", "=Assumptions!C7", "=Assumptions!D7", "Audience-to-member conversion"),
        ("Members", "=ROUND(B2*B3,0)", "=ROUND(C2*C3,0)", "=ROUND(D2*D3,0)", "Derived member base"),
        ("Price / month", "=Assumptions!B8", "=Assumptions!C8", "=Assumptions!D8", "Monthly Passport Club price"),
        ("Months billed", 12, 12, 12, "Annualized view"),
        ("Monthly recurring revenue", "=B3*0", "=C3*0", "=D3*0", "Spacer row for styling"),
        ("Annual membership revenue", "=B4*B5*B6", "=C4*C5*C6", "=D4*D5*D6", "Matches base planning case"),
    ]
    for row in membership_rows:
        ws.append(list(row))
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=8):
        label = row[0].value
        for idx, cell in enumerate(row, start=1):
            style_cell(cell, bold=label == "Annual membership revenue", fill=green_fill if label == "Annual membership revenue" else None)
            if 2 <= idx <= 4:
                if label == "Conversion rate":
                    cell.number_format = "0.0%"
                elif label in {"Price / month", "Monthly recurring revenue", "Annual membership revenue"}:
                    cell.number_format = "$#,##0"
                else:
                    cell.number_format = "#,##0"
    ws["B7"] = "=B4*B3"
    ws["C7"] = "=C4*C3"
    ws["D7"] = "=D4*D3"
    set_widths(ws, [28, 15, 15, 15, 32])

    ws = wb.create_sheet("Revenue_Build")
    ws.append(["Revenue stream", "Year 1", "Year 2", "Year 3", "Mix note"])
    revenue_build_rows = [
        ("Product revenue", "=Assumptions!B6", "=Assumptions!C6", "=Assumptions!D6", "Destination-led collection sales"),
        ("Membership revenue", "=Membership_Model!B8", "=Membership_Model!C8", "=Membership_Model!D8", "Recurring customer layer"),
        ("Partnership revenue", "=Assumptions!B10", "=Assumptions!C10", "=Assumptions!D10", "Affiliate and partnership income"),
        ("Total revenue", "=SUM(B2:B4)", "=SUM(C2:C4)", "=SUM(D2:D4)", "Top-line revenue view"),
        ("Product mix %", "=B2/B5", "=C2/C5", "=D2/D5", "Share of total revenue"),
        ("Membership mix %", "=B3/B5", "=C3/C5", "=D3/D5", "Share of total revenue"),
        ("Partnership mix %", "=B4/B5", "=C4/C5", "=D4/D5", "Share of total revenue"),
        ("Year-over-year growth", "", "=C5/B5-1", "=D5/C5-1", "Total revenue growth"),
    ]
    for row in revenue_build_rows:
        ws.append(list(row))
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=9):
        label = row[0].value
        is_total = label == "Total revenue"
        for idx, cell in enumerate(row, start=1):
            style_cell(cell, bold=is_total, fill=green_fill if is_total else None)
            if 2 <= idx <= 4:
                if "%" in str(label) or label == "Year-over-year growth":
                    cell.number_format = "0.0%"
                elif "revenue" in str(label).lower():
                    cell.number_format = "$#,##0"
    set_widths(ws, [24, 15, 15, 15, 34])
    line_chart = LineChart()
    line_chart.title = "Revenue progression"
    line_chart.y_axis.title = "USD"
    line_chart.style = 13
    data = Reference(ws, min_col=2, max_col=4, min_row=2, max_row=4)
    cats = Reference(ws, min_col=2, max_col=4, min_row=1, max_row=1)
    line_chart.add_data(data, titles_from_data=False, from_rows=True)
    line_chart.set_categories(cats)
    line_chart.height = 7
    line_chart.width = 12
    ws.add_chart(line_chart, "G2")

    ws = wb.create_sheet("Opex_Plan")
    ws.append(["Operating cost", "Year 1", "Year 2", "Year 3", "Notes"])
    opex_rows = [
        ("Travel & sourcing", "=Assumptions!B13", "=Assumptions!C13", "=Assumptions!D13", "Trips, supplier onboarding, and quality control"),
        ("Team", "=Assumptions!B14", "=Assumptions!C14", "=Assumptions!D14", "Core staffing and specialist support"),
        ("Marketing & growth", "=Assumptions!B15", "=Assumptions!C15", "=Assumptions!D15", "Content amplification and demand generation"),
        ("General operations", "=Assumptions!B16", "=Assumptions!C16", "=Assumptions!D16", "Platform, admin, and overhead"),
        ("Total operating costs", "=SUM(B2:B5)", "=SUM(C2:C5)", "=SUM(D2:D5)", "Total operating cost base"),
        ("Opex as % of revenue", "=B6/Revenue_Build!B5", "=C6/Revenue_Build!C5", "=D6/Revenue_Build!D5", "Operating cost ratio"),
    ]
    for row in opex_rows:
        ws.append(list(row))
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=7):
        label = row[0].value
        is_total = label == "Total operating costs"
        for idx, cell in enumerate(row, start=1):
            style_cell(cell, bold=is_total, fill=green_fill if is_total else None)
            if 2 <= idx <= 4:
                if label == "Opex as % of revenue":
                    cell.number_format = "0.0%"
                else:
                    cell.number_format = "$#,##0"
    set_widths(ws, [24, 15, 15, 15, 38])

    ws = wb.create_sheet("Profitability")
    ws.append(["Metric", "Year 1", "Year 2", "Year 3"])
    profitability_rows = [
        ("Product revenue", "=Revenue_Build!B2", "=Revenue_Build!C2", "=Revenue_Build!D2"),
        ("Membership revenue", "=Revenue_Build!B3", "=Revenue_Build!C3", "=Revenue_Build!D3"),
        ("Partnership revenue", "=Revenue_Build!B4", "=Revenue_Build!C4", "=Revenue_Build!D4"),
        ("Total revenue", "=Revenue_Build!B5", "=Revenue_Build!C5", "=Revenue_Build!D5"),
        ("COGS", "=Assumptions!B12", "=Assumptions!C12", "=Assumptions!D12"),
        ("Gross profit", "=B5-B6", "=C5-C6", "=D5-D6"),
        ("Gross margin %", "=B7/B5", "=C7/C5", "=D7/D5"),
        ("Total operating costs", "=Opex_Plan!B6", "=Opex_Plan!C6", "=Opex_Plan!D6"),
        ("Operating profit", "=B7-B9", "=C7-C9", "=D7-D9"),
        ("Operating margin %", "=B10/B5", "=C10/C5", "=D10/D5"),
    ]
    for row in profitability_rows:
        ws.append(list(row))
    for cell in ws[1]:
        style_header(cell)
    highlighted = {"Total revenue", "Gross profit", "Operating profit"}
    percent_labels = {"Gross margin %", "Operating margin %"}
    for row in ws.iter_rows(min_row=2, max_row=11):
        label = row[0].value
        for idx, cell in enumerate(row, start=1):
            style_cell(cell, bold=label in highlighted, fill=green_fill if label in highlighted else None)
            if 2 <= idx <= 4:
                if label in percent_labels:
                    cell.number_format = "0.0%"
                else:
                    cell.number_format = "$#,##0"
    set_widths(ws, [24, 15, 15, 15])
    profit_chart = BarChart()
    profit_chart.type = "col"
    profit_chart.style = 10
    profit_chart.title = "Profitability progression"
    profit_chart.y_axis.title = "USD"
    profit_data = Reference(ws, min_col=2, max_col=4, min_row=5, max_row=7)
    profit_chart.add_data(profit_data, titles_from_data=False, from_rows=True)
    profit_chart.set_categories(Reference(ws, min_col=2, max_col=4, min_row=1, max_row=1))
    profit_chart.height = 7
    profit_chart.width = 12
    ws.add_chart(profit_chart, "F2")

    ws = wb.create_sheet("Capital_Plan")
    ws.append(["Use of funds", "Initial capital", "Growth capital", "Notes"])
    for category, initial_capital, growth_capital, note in CAPITAL_PLAN:
        ws.append([category, initial_capital, growth_capital, note])
    ws.append(["Total", "=SUM(B2:B7)", "=SUM(C2:C7)", "Illustrative capital plan"])
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=7):
        for idx, cell in enumerate(row, start=1):
            style_cell(cell)
            if idx in {2, 3}:
                cell.number_format = "$#,##0"
    for cell in ws[8]:
        style_cell(cell, bold=True, fill=green_fill)
        if cell.column in {2, 3}:
            cell.number_format = "$#,##0"
    set_widths(ws, [30, 16, 16, 46])

    ws = wb.create_sheet("Scenario_View")
    ws.append(["Scenario", "Revenue multiplier", "Year 3 revenue", "Operating margin", "Year 3 operating profit", "Commentary"])
    scenario_rows = [
        ("Downside", 0.70, "=Profitability!D5*B2", "=Profitability!D11*0.78", "=C2*D2", "Slower audience growth and softer sell-through compress the ramp."),
        ("Base case", 1.00, "=Profitability!D5*B3", "=Profitability!D11", "=C3*D3", "Current modeled trajectory."),
        ("Upside", 1.28, "=Profitability!D5*B4", "=Profitability!D11*1.08", "=C4*D4", "Faster repeats, stronger private label traction, and cleaner retention."),
    ]
    for row in scenario_rows:
        ws.append(list(row))
    for cell in ws[1]:
        style_header(cell)
    for row in ws.iter_rows(min_row=2, max_row=4):
        for idx, cell in enumerate(row, start=1):
            style_cell(cell)
            if idx == 2 or idx == 4:
                cell.number_format = "0.0%"
            elif idx in {3, 5}:
                cell.number_format = "$#,##0"
    set_widths(ws, [16, 18, 18, 16, 20, 56])

    for ws in wb.worksheets:
        ws.freeze_panes = "A2"
        ws.sheet_view.showGridLines = False

    wb.save(str(output))


def build_zip() -> None:
    output = DOCS / "alidandin_investor_package.zip"
    with ZipFile(output, "w", compression=ZIP_DEFLATED) as archive:
        for filename in ["investor_business_plan.pdf", "investor_deck.pptx", "financial_model.xlsx"]:
            archive.write(DOCS / filename, arcname=filename)


def main() -> None:
    ensure_dirs()
    build_pdf()
    build_pptx()
    build_xlsx()
    build_zip()


if __name__ == "__main__":
    main()
